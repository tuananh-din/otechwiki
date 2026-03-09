import base64
import fitz  # PyMuPDF
import hashlib
import httpx
import io
import pptx
import tiktoken
from bs4 import BeautifulSoup
from markitdown import MarkItDown
from openai import AsyncOpenAI
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from app.models.document import Document, Chunk
from app.services.embeddings import get_embeddings_batch
from app.services.url_utils import normalize_url
from app.core.config import get_settings
from datetime import datetime, timezone

settings = get_settings()
tokenizer = tiktoken.encoding_for_model("gpt-4o-mini")
openai_client = AsyncOpenAI(api_key=settings.openai_api_key)

CHUNK_SIZE = 800  # tokens (slightly larger for better context)
CHUNK_OVERLAP = 150


def split_into_chunks(text: str, document_title: str = "", chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[dict]:
    """Split text into overlapping chunks and prepend document context."""
    tokens = tokenizer.encode(text)
    chunks = []
    start = 0
    idx = 0

    while start < len(tokens):
        end = min(start + chunk_size, len(tokens))
        chunk_tokens = tokens[start:end]
        chunk_text = tokenizer.decode(chunk_tokens)
        
        # Prepend context to improve retrieval
        contextual_text = f"Tài liệu: {document_title}\n---\n{chunk_text}"
        
        chunks.append({
            "text": contextual_text, 
            "raw_content": chunk_text,
            "index": idx, 
            "token_count": len(chunk_tokens)
        })
        idx += 1
        start += chunk_size - overlap
        if end >= len(tokens):
            break

    return chunks


async def extract_text_with_vision(image_bytes: bytes) -> str:
    """Use GPT-4o-mini Vision to extract text and tables from an image in Markdown format."""
    base64_image = base64.b64encode(image_bytes).decode("utf-8")
    
    response = await openai_client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "system",
                "content": "You are a specialized document parser. Extract all text, tables, and lists from this document page image. Output ONLY the extracted content in clean Markdown format. Preserve the visual structure (tables as | | tables, lists as - bullets). If the page is in Vietnamese, extract as Vietnamese."
            },
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": "Extract all content from this page image into Markdown."
                    },
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/png;base64,{base64_image}"}
                    }
                ]
            }
        ],
        max_tokens=2048
    )
    return response.choices[0].message.content or ""


def recursive_extract_pptx_text(shape) -> str:
    """Recursively extract text from PPTX shapes, including groups and tables."""
    text_content = []
    
    # Process text frames
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            p_text = paragraph.text.strip()
            if p_text:
                text_content.append(p_text)
    
    # Process tables
    if hasattr(shape, "has_table") and shape.has_table:
        for row in shape.table.rows:
            row_text = " | ".join([cell.text_frame.text.strip() for cell in row.cells])
            if row_text.strip():
                text_content.append(row_text)
                
    # Process groups (ShapeType 6 is GROUP)
    if hasattr(shape, "shape_type") and shape.shape_type == 6:
        for subshape in shape.shapes:
            text_content.append(recursive_extract_pptx_text(subshape))
            
    return "\n".join(text_content)


async def ingest_pdf(db: AsyncSession, document_id: int, file_path: str) -> int:
    """Extract text using MarkItDown + Vision Fallback for high-fidelity RAG."""
    doc_record = await db.get(Document, document_id)
    if not doc_record:
        raise ValueError(f"Document {document_id} not found")

    doc_record.status = "processing"
    await db.commit()

    try:
        # 1. Fast Extraction with MarkItDown
        md = MarkItDown()
        result = md.convert(file_path)
        markdown_text = result.text_content
        
        # 2. Heuristic: If extraction is too light (< 200 chars/page), use Vision
        pdf = fitz.open(file_path)
        num_pages = len(pdf)
        doc_record.page_count = num_pages
        
        # Calculate char density (avoid division by zero)
        density = len(markdown_text) / max(num_pages, 1)
        
        if density < 150:  # Threshold for "image-heavy" or "complex layout"
            print(f"Low density ({density}) detected. Switching to Vision-OCR for {document_id}")
            vision_parts = []
            for page_num in range(num_pages):
                page = pdf[page_num]
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2)) # Zoom for better OCR
                img_bytes = pix.tobytes("png")
                page_md = await extract_text_with_vision(img_bytes)
                vision_parts.append(f"<!-- Page {page_num + 1} -->\n{page_md}")
            markdown_text = "\n\n---\n\n".join(vision_parts)
        
        pdf.close()
        
        if not markdown_text.strip():
            doc_record.status = "error"
            await db.commit()
            return 0

        all_chunks = split_into_chunks(markdown_text, document_title=doc_record.title)

        # Generate embeddings and store
        texts = [c["text"] for c in all_chunks]
        embeddings = await get_embeddings_batch(texts)

        for chunk, embedding in zip(all_chunks, embeddings):
            db_chunk = Chunk(
                document_id=document_id,
                content=chunk["text"],
                embedding=embedding,
                chunk_index=chunk["index"],
                token_count=chunk["token_count"],
            )
            db.add(db_chunk)

        doc_record.status = "ready"
        await db.commit()
        return len(all_chunks)

    except Exception as e:
        doc_record.status = "error"
        await db.commit()
        raise e


async def ingest_web(db: AsyncSession, document_id: int, url: str) -> int:
    """V2 Pipeline: fetch → clean → dedup → smart_chunk → embed → map."""
    from app.services.cleaner import clean_html, detect_page_type, extract_domain
    from app.services.dedup import dedup_blocks, dedup_chunks
    from app.services.chunker_v2 import smart_chunk
    from app.services.mapper_v2 import map_document_v2, apply_mappings

    doc_record = await db.get(Document, document_id)
    if not doc_record:
        raise ValueError(f"Document {document_id} not found")

    doc_record.status = "processing"
    await db.commit()

    try:
        # 1. Fetch HTML
        headers = {"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"}
        async with httpx.AsyncClient(timeout=30, follow_redirects=True, headers=headers) as client:
            resp = await client.get(url)
            resp.raise_for_status()
            raw_html = resp.text

        # 1b. Hash-based change detection
        new_raw_hash = hashlib.sha256(raw_html.encode()).hexdigest()
        doc_record.last_fetched_at = datetime.now(timezone.utc)
        doc_record.canonical_url = normalize_url(url)

        # Save HTTP caching headers
        doc_record.etag = resp.headers.get("ETag", "")
        doc_record.last_modified_header = resp.headers.get("Last-Modified", "")

        if doc_record.raw_hash == new_raw_hash:
            doc_record.import_status = "unchanged"
            doc_record.freshness_score = 100
            await db.commit()
            return 0  # Skip — content identical

        doc_record.raw_hash = new_raw_hash

        # 2. Detect page type and domain
        page_type = detect_page_type(url, raw_html)
        domain = extract_domain(url)
        doc_record.page_type = page_type
        doc_record.domain = domain

        # 3. Clean HTML (rule-based, 0 tokens)
        if settings.web_extractor_url:
            # Use Defuddle worker for extraction
            async with httpx.AsyncClient(timeout=30, follow_redirects=True, headers=headers) as client:
                extractor_url = settings.web_extractor_url.rstrip("/")
                fetch_url = f"{extractor_url}/{url}"
                ext_resp = await client.get(fetch_url)
                ext_resp.raise_for_status()
                raw_text = ext_resp.text
                cleaned_text = raw_text  # Defuddle already cleans
        else:
            # BS4 + custom boilerplate removal
            clean_result = clean_html(raw_html, page_type)
            raw_text = clean_result["raw_text"]
            cleaned_text = clean_result["cleaned_text"]

        if not cleaned_text or not cleaned_text.strip():
            doc_record.status = "error"
            doc_record.cleaning_status = "error"
            await db.commit()
            return 0

        # 4. Save raw + cleaned text + clean hash
        new_clean_hash = hashlib.sha256(cleaned_text.encode()).hexdigest()

        if doc_record.clean_hash == new_clean_hash:
            # Raw HTML changed (template/UI update) but actual content identical
            doc_record.import_status = "unchanged"
            doc_record.freshness_score = 100
            await db.commit()
            return 0  # Skip re-chunking

        doc_record.clean_hash = new_clean_hash
        doc_record.raw_text = raw_text
        doc_record.cleaned_text = cleaned_text
        doc_record.cleaning_status = "cleaned"
        doc_record.import_status = "cleaned"

        # 5. Dedup blocks within page
        deduped_text = dedup_blocks(cleaned_text)

        # 6. Smart chunk (heading-aware)
        all_chunks = smart_chunk(deduped_text, document_title=doc_record.title)

        if not all_chunks:
            doc_record.status = "error"
            await db.commit()
            return 0

        # 7. Dedup chunks
        all_chunks = dedup_chunks(all_chunks)

        # 8. Embed only searchable chunks
        searchable = [c for c in all_chunks if c.get("is_searchable", True)]
        if searchable:
            texts = [c["cleaned_content"] for c in searchable]
            embeddings = await get_embeddings_batch(texts)
        else:
            embeddings = []

        # 9. Delete old chunks for this document
        from sqlalchemy import delete as sa_delete
        await db.execute(sa_delete(Chunk).where(Chunk.document_id == document_id))

        # 10. Save all chunks (searchable gets embedding, non-searchable gets null)
        embed_idx = 0
        for chunk_data in all_chunks:
            is_search = chunk_data.get("is_searchable", True)
            db_chunk = Chunk(
                document_id=document_id,
                content=chunk_data["content"],
                cleaned_content=chunk_data.get("cleaned_content"),
                embedding=embeddings[embed_idx] if is_search and embed_idx < len(embeddings) else None,
                chunk_index=chunk_data["chunk_index"],
                token_count=chunk_data.get("token_count", 0),
                section_title=chunk_data.get("section_title"),
                section_path=chunk_data.get("section_path"),
                is_searchable=is_search,
                dedup_hash=chunk_data.get("dedup_hash"),
            )
            db.add(db_chunk)
            if is_search:
                embed_idx += 1

        # 11. Confidence mapping
        try:
            mappings = await map_document_v2(db, doc_record)
            if mappings:
                await apply_mappings(db, doc_record, mappings)
        except Exception:
            pass  # Mapping errors should not block ingestion

        doc_record.status = "ready"
        doc_record.import_status = "indexed"
        doc_record.freshness_score = 100
        await db.commit()
        return len(all_chunks)

    except Exception as e:
        doc_record.status = "error"
        doc_record.cleaning_status = "error"
        await db.commit()
        raise e
async def ingest_ppt(db: AsyncSession, document_id: int, file_path: str) -> int:
    """Extract text from PPTX using Recursive Shape Analysis + MarkItDown fallback."""
    doc_record = await db.get(Document, document_id)
    if not doc_record:
        raise ValueError(f"Document {document_id} not found")

    doc_record.status = "processing"
    await db.commit()

    try:
        prs = pptx.Presentation(file_path)
        doc_record.page_count = len(prs.slides)
        
        slide_parts = []
        for i, slide in enumerate(prs.slides):
            page_text = f"<!-- Slide {i+1} -->\n"
            shapes_text = []
            for shape in slide.shapes:
                shapes_text.append(recursive_extract_pptx_text(shape))
            
            # If shape extraction is very light, use MarkItDown for the whole slide?
            # Or just trust the recursive walker.
            slide_md = "\n".join([t for t in shapes_text if t.strip()])
            slide_parts.append(page_text + slide_md)
            
        markdown_text = "\n\n---\n\n".join(slide_parts)

        if not markdown_text.strip():
            # Fallback to MarkItDown
            md = MarkItDown()
            result = md.convert(file_path)
            markdown_text = result.text_content

        if not markdown_text.strip():
            doc_record.status = "error"
            await db.commit()
            return 0

        all_chunks = split_into_chunks(markdown_text, document_title=doc_record.title)

        # Generate embeddings and store
        texts = [c["text"] for c in all_chunks]
        embeddings = await get_embeddings_batch(texts)

        for chunk, embedding in zip(all_chunks, embeddings):
            db_chunk = Chunk(
                document_id=document_id,
                content=chunk["text"],
                embedding=embedding,
                chunk_index=chunk["index"],
                token_count=chunk["token_count"],
            )
            db.add(db_chunk)

        doc_record.status = "ready"
        await db.commit()
        return len(all_chunks)

    except Exception as e:
        doc_record.status = "error"
        await db.commit()
        raise e
