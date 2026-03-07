import asyncio
import os
import sys
import fitz
import pptx
import tiktoken
from typing import List, Dict

# Add backend to sys.path
sys.path.append(os.path.join(os.getcwd(), "backend"))

from app.services.ingest import split_into_chunks

def audit_pdf(file_path: str):
    print(f"\n--- Auditing PDF: {os.path.basename(file_path)} ---")
    try:
        doc = fitz.open(file_path)
        total_pages = len(doc)
        print(f"Total Pages: {total_pages}")
        
        full_text = ""
        for i in range(total_pages):
            page = doc[i]
            text = page.get_text().strip()
            print(f"Page {i+1}: {len(text)} characters extracted.")
            full_text += text + "\n"
        
        doc.close()
        
        chunks = split_into_chunks(full_text)
        print(f"Total Chunks: {len(chunks)}")
        print(f"Avg Chunk Size: {sum(c['token_count'] for c in chunks)/len(chunks) if chunks else 0:.1f} tokens")
        
    except Exception as e:
        print(f"Error auditing PDF: {e}")

def audit_ppt(file_path: str):
    print(f"\n--- Auditing PPTX: {os.path.basename(file_path)} ---")
    try:
        prs = pptx.Presentation(file_path)
        total_slides = len(prs.slides)
        print(f"Total Slides: {total_slides}")
        
        full_text = ""
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_text += paragraph.text.strip() + " "
                elif hasattr(shape, "text"):
                    slide_text += shape.text.strip() + " "
            
            print(f"Slide {i+1}: {len(slide_text.strip())} characters extracted.")
            full_text += slide_text.strip() + "\n"
            
        chunks = split_into_chunks(full_text)
        print(f"Total Chunks: {len(chunks)}")
        print(f"Avg Chunk Size: {sum(c['token_count'] for c in chunks)/len(chunks) if chunks else 0:.1f} tokens")
        
    except Exception as e:
        print(f"Error auditing PPTX: {e}")

if __name__ == "__main__":
    uploads_dir = "backend/uploads"
    if not os.path.exists(uploads_dir):
        print(f"Uploads directory not found: {uploads_dir}")
        sys.exit(1)
        
    files = os.listdir(uploads_dir)
    for f in files:
        path = os.path.join(uploads_dir, f)
        if f.lower().endswith(".pdf"):
            audit_pdf(path)
        elif f.lower().endswith(".pptx"):
            audit_ppt(path)
